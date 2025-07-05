import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt # Import Pt for font sizes
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import math # Import math for ceil function

# Get the logger from the excel_loader module to ensure consistent logging
# This assumes excel_loader.py is configured and its logger is accessible
try:
    from excel_loader import logger
except ImportError:
    # Fallback logger configuration if excel_loader is not available
    logging.basicConfig(level=logging.INFO,
                        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                        handlers=[logging.FileHandler("pptx_exporter.log")]) # Removed StreamHandler
    logger = logging.getLogger(__name__)
    logger.warning("Could not import logger from excel_loader. Using fallback logger configuration (file only).")


def export_dataframe_to_pptx(dataframe, pptx_file_name, slide_title="Data Table", rows_per_slide=5):
    """
    Exports a pandas DataFrame to a new PowerPoint presentation as a table,
    paginating the data across multiple slides.

    Args:
        dataframe (pandas.DataFrame): The DataFrame to export.
        pptx_file_name (str): The name of the PowerPoint file to create/overwrite.
        slide_title (str): The base title for the slides containing the table.
        rows_per_slide (int): The number of DataFrame rows to display per slide.
    """
    if dataframe is None or dataframe.empty:
        logger.warning("DataFrame is empty or None. Nothing to export to PowerPoint.")
        return

    try:
        # Create a new presentation
        prs = Presentation()
        # Choose a slide layout (e.g., Title and Content layout)
        # Layout 1 is typically 'Title and Content'
        slide_layout = prs.slide_layouts[1]

        total_rows = len(dataframe)
        num_slides = math.ceil(total_rows / rows_per_slide)

        logger.info(f"Exporting {total_rows} rows across {num_slides} slides, with {rows_per_slide} rows per slide.")

        for i in range(num_slides):
            start_row = i * rows_per_slide
            end_row = min((i + 1) * rows_per_slide, total_rows)
            display_dataframe = dataframe.iloc[start_row:end_row]

            if display_dataframe.empty:
                continue # Skip if a slice somehow ends up empty

            slide = prs.slides.add_slide(slide_layout)

            # Set the slide title
            title = slide.shapes.title
            title.text = f"{slide_title} (Page {i+1})"
            logger.info(f"Added slide {i+1}/{num_slides} with title: '{title.text}'")

            # Add a table to the slide
            rows_in_current_chunk, cols = display_dataframe.shape
            # Add 1 to rows for the header row
            table_rows = rows_in_current_chunk + 1
            table_cols = cols

            # Define table position and size (adjust as needed)
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)

            # Add the table shape
            shape = slide.shapes.add_table(table_rows, table_cols, left, top, width, height)
            table = shape.table

            # Set column widths (optional, but good for readability)
            # Distribute width equally among columns
            for col_idx in range(table_cols):
                # Convert the calculated width to an integer
                table.columns[col_idx].width = int(width / table_cols)

            # Populate header row
            for col_idx, col_name in enumerate(display_dataframe.columns):
                cell = table.cell(0, col_idx)
                text_frame = cell.text_frame
                text_frame.text = str(col_name)
                # Basic header styling
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].font.size = Pt(14) # Adjust font size
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 112, 192) # Dark Blue
                text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # White text
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


            # Populate data rows
            for row_idx, row_data in enumerate(display_dataframe.itertuples(index=False)): # Use itertuples for efficiency
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx) # +1 because of header row
                    text_frame = cell.text_frame
                    text_frame.text = str(cell_value)
                    text_frame.paragraphs[0].font.size = Pt(11) # Adjust font size
                    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # Align text left

            # Add page number to the lower right-hand side
            # Define position for the page number
            page_num_left = Inches(8.5)
            page_num_top = Inches(7.0)
            page_num_width = Inches(1.0)
            page_num_height = Inches(0.5)

            # Add a text box for the page number
            page_num_shape = slide.shapes.add_textbox(page_num_left, page_num_top, page_num_width, page_num_height)
            page_num_text_frame = page_num_shape.text_frame
            page_num_text_frame.text = f"Page {i+1} of {num_slides}"
            page_num_text_frame.paragraphs[0].font.size = Pt(10)
            page_num_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            page_num_text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128) # Grey color
            page_num_shape.fill.background() # Make background transparent
            page_num_shape.line.fill.background() # Make border transparent
            logger.info(f"Added page number {i+1} to slide {i+1}.")


        # Save the presentation
        prs.save(pptx_file_name)
        logger.info(f"Successfully exported DataFrame to '{pptx_file_name}' with {num_slides} slides.")

    except Exception as e:
        logger.error(f"An error occurred while exporting DataFrame to PowerPoint: {e}", exc_info=True)
